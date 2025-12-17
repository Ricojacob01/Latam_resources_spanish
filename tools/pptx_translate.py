import json
import os
import sys
from typing import List, Dict, Any

from pptx import Presentation


def _iter_shape_texts(shape, texts: List[Dict[str, Any]], slide_index: int):
    # Text frames
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        content_lines = []
        try:
            for p in shape.text_frame.paragraphs:
                runs = [r.text for r in p.runs] or [p.text]
                content_lines.append("".join(runs))
        except Exception:
            try:
                content_lines.append(shape.text)
            except Exception:
                pass
        full_text = "\n".join([t for t in content_lines if t is not None])
        if full_text.strip():
            texts.append({
                "slide_index": slide_index,
                "shape_id": getattr(shape, "shape_id", None),
                "type": "text_frame",
                "text": full_text
            })
    # Tables
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for r, row in enumerate(table.rows):
            for c, cell in enumerate(row.cells):
                cell_text = cell.text or ""
                if cell_text.strip():
                    texts.append({
                        "slide_index": slide_index,
                        "shape_id": getattr(shape, "shape_id", None),
                        "type": "table_cell",
                        "row": r,
                        "col": c,
                        "text": cell_text
                    })
    # Groups
    if hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            _iter_shape_texts(subshape, texts, slide_index)


def extract_texts(pptx_path: str, out_json: str):
    prs = Presentation(pptx_path)
    all_texts: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            _iter_shape_texts(shape, all_texts, idx)
    with open(out_json, "w", encoding="utf-8") as f:
        json.dump(all_texts, f, ensure_ascii=False, indent=2)
    print(f"Extracted {len(all_texts)} text entries to {out_json}")


def apply_translations(pptx_path: str, translations_json: str, out_pptx: str):
    with open(translations_json, "r", encoding="utf-8") as f:
        translations = json.load(f)
    prs = Presentation(pptx_path)
    # Build index by (slide_index, shape_id, type, row, col)
    trans_index = {}
    for t in translations:
        key = (
            t.get("slide_index"),
            t.get("shape_id"),
            t.get("type"),
            t.get("row"),
            t.get("col"),
        )
        trans_index[key] = t.get("translated_text", t.get("text"))
    updated = 0
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Text frames
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                key = (idx, getattr(shape, "shape_id", None), "text_frame", None, None)
                if key in trans_index:
                    # Replace entire text content
                    new_text = trans_index[key] or ""
                    try:
                        # Clear existing paragraphs
                        tf = shape.text_frame
                        while len(tf.paragraphs) > 1:
                            p = tf.paragraphs[-1]
                            p._p.getparent().remove(p._p)  # remove extra paragraphs
                        tf.clear()
                        tf.text = new_text
                        updated += 1
                    except Exception:
                        pass
            # Tables
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for r, row in enumerate(table.rows):
                    for c, cell in enumerate(row.cells):
                        key = (idx, getattr(shape, "shape_id", None), "table_cell", r, c)
                        if key in trans_index:
                            try:
                                cell.text = trans_index[key] or ""
                                updated += 1
                            except Exception:
                                pass
            # Groups
            if hasattr(shape, "shapes"):
                for subshape in shape.shapes:
                    # Only simple handling: subshapes handled via iteration in extraction aren't directly addressed here
                    # because shape_id mapping can differ; rely on parent shape_id match for text_frame replacements
                    pass
    prs.save(out_pptx)
    print(f"Applied translations to {updated} entries. Saved: {out_pptx}")


def main():
    if len(sys.argv) < 3:
        print("Usage:")
        print("  Extract: python tools/pptx_translate.py extract <in.pptx> <out.json>")
        print("  Apply  : python tools/pptx_translate.py apply <in.pptx> <translations.json> <out.pptx>")
        sys.exit(1)
    cmd = sys.argv[1]
    if cmd == "extract":
        _, _, in_pptx, out_json = sys.argv
        extract_texts(in_pptx, out_json)
    elif cmd == "apply":
        _, _, in_pptx, translations_json, out_pptx = sys.argv
        apply_translations(in_pptx, translations_json, out_pptx)
    else:
        print(f"Unknown command: {cmd}")
        sys.exit(2)


if __name__ == "__main__":
    main()


